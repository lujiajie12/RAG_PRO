from __future__ import annotations

from io import BytesIO

from docx import Document as DocxDocument
from pptx import Presentation
from pypdf import PdfWriter

from app.ingestion.loaders.csv_loader import CsvLoader
from app.ingestion.loaders.docx_loader import DocxLoader
from app.ingestion.loaders.html_loader import HtmlLoader
from app.ingestion.loaders.md_loader import MarkdownLoader
from app.ingestion.loaders.pdf_loader import PdfLoader
from app.ingestion.loaders.ppt_loader import PptLoader
from app.ingestion.parsers import ParserRegistry
from app.ingestion.types import ParsedDocument, ParsedSegment
from app.rag.cleaner import TextCleaner
from app.rag.loaders import DOCUMENT_LOADER_MAPPING, DocumentLoaderFactory
from app.rag.metadata import DocumentMetadataExtractor
from app.rag.splitter import ParentChildSplitter


def test_rag_loader_mapping_exposes_supported_suffixes():
    assert DOCUMENT_LOADER_MAPPING[".pdf"] == "PdfLoader"
    assert DOCUMENT_LOADER_MAPPING[".docx"] == "DocxLoader"
    assert DOCUMENT_LOADER_MAPPING[".pptx"] == "PptLoader"
    assert DOCUMENT_LOADER_MAPPING[".csv"] == "CsvLoader"


def test_loader_parser_cleaner_metadata_splitter_smoke(tmp_path):
    path = tmp_path / "guide.md"
    path.write_text("# Title\n\nalpha beta gamma\ndelta epsilon zeta\n", encoding="utf-8")

    loaded = DocumentLoaderFactory().load_document(path)
    parsed = ParserRegistry().parse(loaded)
    cleaned = TextCleaner().clean(parsed)
    enriched = DocumentMetadataExtractor().extract(loaded, parsed, cleaned)
    chunks = ParentChildSplitter(parent_tokens=12, child_tokens=6, child_overlap_tokens=3).split(enriched.cleaned)

    assert loaded.file_name == "guide.md"
    assert parsed.parser_name == "markdown"
    assert cleaned.segments[1].cleaned_text == "alpha beta gamma delta epsilon zeta"
    assert enriched.metadata["title"] == "Title"
    assert enriched.cleaned.segments[0].metadata["section_path"] == ["Title"]
    assert chunks["parents"]
    assert chunks["children"]


def test_text_cleaner_merges_paragraph_lines_but_preserves_structure():
    parsed = ParsedDocument(
        file_name="sample.md",
        file_type="md",
        parser_name="markdown",
        segments=[
            ParsedSegment(kind="heading", text="Title", raw_text="Title", metadata={"heading_level": 1}),
            ParsedSegment(kind="paragraph", text="Line one\nLine two", raw_text="Line one\nLine two"),
            ParsedSegment(kind="table_row", text="A | B", raw_text="A | B"),
        ],
    )

    cleaned = TextCleaner().clean(parsed)

    assert cleaned.segments[0].cleaned_text == "Title"
    assert cleaned.segments[1].cleaned_text == "Line one Line two"
    assert cleaned.segments[2].cleaned_text == "A | B"
    assert cleaned.cleaning_stats["merged_lines"] == 1


def test_text_cleaner_drops_repeated_pdf_headers_and_footers():
    parsed = ParsedDocument(
        file_name="report.pdf",
        file_type="pdf",
        parser_name="pdf",
        metadata={"page_count": 3},
        segments=[
            ParsedSegment(kind="page", text="ACME\nBody 1\nConfidential", raw_text="ACME\nBody 1\nConfidential", metadata={"page_number": 1}),
            ParsedSegment(kind="page", text="ACME\nBody 2\nConfidential", raw_text="ACME\nBody 2\nConfidential", metadata={"page_number": 2}),
            ParsedSegment(kind="page", text="ACME\nBody 3\nConfidential", raw_text="ACME\nBody 3\nConfidential", metadata={"page_number": 3}),
        ],
    )

    cleaned = TextCleaner().clean(parsed)

    assert [segment.cleaned_text for segment in cleaned.segments] == ["Body 1", "Body 2", "Body 3"]
    assert cleaned.cleaning_stats["dropped_repeated_headers"] == 3
    assert cleaned.cleaning_stats["dropped_repeated_footers"] == 3


def test_metadata_extractor_adds_section_path_and_token_counts(tmp_path):
    path = tmp_path / "outline.md"
    path.write_text("# Root\n\n## Child\n\nParagraph body\n", encoding="utf-8")

    loaded = DocumentLoaderFactory().load_document(path)
    parsed = ParserRegistry().parse(loaded)
    cleaned = TextCleaner().clean(parsed)
    enriched = DocumentMetadataExtractor().extract(loaded, parsed, cleaned)

    assert enriched.metadata["token_count"] > 0
    assert enriched.cleaned.segments[0].metadata["section_path"] == ["Root"]
    assert enriched.cleaned.segments[1].metadata["section_path"] == ["Root", "Child"]
    assert enriched.cleaned.segments[2].metadata["section_path"] == ["Root", "Child"]
    assert enriched.cleaned.segments[2].metadata["token_count"] > 0


def test_parent_child_splitter_uses_overlap_and_inherits_section_path(tmp_path):
    path = tmp_path / "chunks.md"
    path.write_text(
        "# Title\n\nalpha beta gamma\n\ndelta epsilon zeta\n\neta theta iota\n",
        encoding="utf-8",
    )

    loaded = DocumentLoaderFactory().load_document(path)
    parsed = MarkdownLoader().parse_bytes(path.read_bytes(), path.name)
    cleaned = TextCleaner().clean(parsed)
    enriched = DocumentMetadataExtractor().extract(loaded, parsed, cleaned)
    child_token_budget = (
        enriched.cleaned.segments[1].metadata["token_count"] + enriched.cleaned.segments[2].metadata["token_count"]
    )
    overlap_budget = enriched.cleaned.segments[2].metadata["token_count"]

    chunks = ParentChildSplitter(
        parent_tokens=40,
        child_tokens=child_token_budget,
        child_overlap_tokens=overlap_budget,
    ).split(enriched.cleaned)

    assert len(chunks["children"]) >= 2
    overlapping_children = [chunk for chunk in chunks["children"] if "delta epsilon zeta" in chunk["content"]]
    assert len(overlapping_children) >= 2
    assert overlapping_children[0]["metadata"]["section_path"] == ["Title"]
    assert overlapping_children[1]["metadata"]["section_path"] == ["Title"]
    assert overlapping_children[0]["metadata"]["token_count"] > 0


def test_html_loader_extracts_visible_text():
    parsed = HtmlLoader().parse_bytes(
        b"<html><body><h1>Title</h1><script>bad()</script><p>Hello world</p></body></html>",
        "page.html",
    )
    content = HtmlLoader().load_bytes(
        b"<html><body><h1>Title</h1><script>bad()</script><p>Hello world</p></body></html>",
        "page.html",
    )

    assert [segment.kind for segment in parsed.segments] == ["heading", "paragraph"]
    assert content == "Title\nHello world"


def test_csv_loader_formats_rows():
    parsed = CsvLoader().parse_bytes("name,age\nAlice,30\nBob,28\n".encode("utf-8"), "people.csv")
    content = CsvLoader().load_bytes("name,age\nAlice,30\nBob,28\n".encode("utf-8"), "people.csv")

    assert parsed.metadata["row_count"] == 2
    assert parsed.metadata["columns"] == ["name", "age"]
    assert "name: Alice | age: 30" in content
    assert "name: Bob | age: 28" in content


def test_docx_loader_reads_paragraphs_and_tables():
    document = DocxDocument()
    document.add_paragraph("Paragraph content")
    table = document.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "A"
    table.cell(0, 1).text = "B"
    buffer = BytesIO()
    document.save(buffer)

    parsed = DocxLoader().parse_bytes(buffer.getvalue(), "sample.docx")
    content = DocxLoader().load_bytes(buffer.getvalue(), "sample.docx")

    assert any(segment.kind == "paragraph" for segment in parsed.segments)
    assert any(segment.kind == "table_row" for segment in parsed.segments)
    assert "Paragraph content" in content
    assert "A | B" in content


def test_ppt_loader_reads_slide_text():
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Intro"
    slide.placeholders[1].text = "First bullet"
    buffer = BytesIO()
    presentation.save(buffer)

    parsed = PptLoader().parse_bytes(buffer.getvalue(), "deck.pptx")
    content = PptLoader().load_bytes(buffer.getvalue(), "deck.pptx")

    assert parsed.metadata["slide_count"] >= 1
    assert any(segment.kind == "slide_title" for segment in parsed.segments)
    assert "Intro" in content
    assert "First bullet" in content


def test_pdf_loader_handles_blank_pdf_without_crashing():
    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    buffer = BytesIO()
    writer.write(buffer)

    parsed = PdfLoader().parse_bytes(buffer.getvalue(), "blank.pdf")
    content = PdfLoader().load_bytes(buffer.getvalue(), "blank.pdf")

    assert parsed.metadata["page_count"] == 1
    assert isinstance(content, str)
