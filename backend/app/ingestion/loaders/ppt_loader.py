from __future__ import annotations

from io import BytesIO

from pptx import Presentation

from .base_loader import BaseLoader
from ..types import ParsedDocument


class PptLoader(BaseLoader):
    def parse_bytes(self, data: bytes, file_name: str) -> ParsedDocument:
        presentation = Presentation(BytesIO(data))
        segments = []

        for slide_index, slide in enumerate(presentation.slides, start=1):
            title = slide.shapes.title.text.strip() if slide.shapes.title and slide.shapes.title.text else ""
            if title:
                segments.append(self.make_segment("slide_title", title, slide_number=slide_index))

            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    for paragraph in shape.text_frame.paragraphs:
                        text = " ".join(paragraph.text.split())
                        if text and text != title:
                            segments.append(self.make_segment("slide_text", text, slide_number=slide_index))
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if cells:
                            segments.append(
                                self.make_segment("table_row", " | ".join(cells), slide_number=slide_index, cell_count=len(cells))
                            )

        return ParsedDocument(
            file_name=file_name,
            file_type="pptx",
            parser_name="pptx",
            segments=segments,
            metadata={"slide_count": len(presentation.slides)},
        )
